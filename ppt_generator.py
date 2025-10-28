import os
import random
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openai_client import generate_content
import json

# Define color themes/templates
TEMPLATES = {
    "professional_blue": {
        "name": "Professional Blue",
        "title_color": RGBColor(0, 51, 102),
        "text_color": RGBColor(51, 51, 51),
        "bg_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(0, 102, 204)
    },
    "modern_green": {
        "name": "Modern Green",
        "title_color": RGBColor(34, 139, 34),
        "text_color": RGBColor(64, 64, 64),
        "bg_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(50, 205, 50)
    },
    "vibrant_orange": {
        "name": "Vibrant Orange",
        "title_color": RGBColor(230, 81, 0),
        "text_color": RGBColor(51, 51, 51),
        "bg_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(255, 140, 0)
    },
    "elegant_purple": {
        "name": "Elegant Purple",
        "title_color": RGBColor(106, 27, 154),
        "text_color": RGBColor(51, 51, 51),
        "bg_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(142, 36, 170)
    },
    "corporate_gray": {
        "name": "Corporate Gray",
        "title_color": RGBColor(64, 64, 64),
        "text_color": RGBColor(96, 96, 96),
        "bg_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(128, 128, 128)
    }
}

def get_template(template_name=None):
    """Get a template by name or return a random one"""
    if template_name and template_name in TEMPLATES:
        return TEMPLATES[template_name]
    return random.choice(list(TEMPLATES.values()))

def create_fallback_slides(prompt: str, context: str):
    """Create fallback slides when AI generation fails"""
    return [
        {
            "title": prompt[:100] if len(prompt) <= 100 else prompt[:97] + "...",
            "content": ["AI-generated presentation", "Based on your requirements"]
        },
        {
            "title": "Overview",
            "content": [
                f"Topic: {prompt[:80]}",
                f"Context: {context[:80]}",
                "This presentation was generated automatically"
            ]
        },
        {
            "title": "Key Points",
            "content": [
                "Point 1: Please review and customize",
                "Point 2: Add your specific content",
                "Point 3: Update as needed"
            ]
        },
        {
            "title": "Summary",
            "content": [
                "This is a template presentation",
                "Please customize with your content",
                "Thank you"
            ]
        }
    ]

def parse_ai_slides(ai_content: str):
    """Parse AI-generated content into structured slides with validation"""
    try:
        slides_data = json.loads(ai_content)
        slides = slides_data.get("slides", [])
        
        if not slides or len(slides) < 2:
            return None
        
        for slide in slides:
            if not slide.get("title") or not isinstance(slide.get("content"), list):
                return None
        
        return slides
    except json.JSONDecodeError:
        lines = ai_content.split('\n')
        slides = []
        current_slide = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('Slide') and ':' in line:
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "title": line.split(':', 1)[1].strip(),
                    "content": []
                }
            elif current_slide and (line.startswith('-') or line.startswith('•')):
                current_slide["content"].append(line.lstrip('-•').strip())
            elif current_slide and line:
                current_slide["content"].append(line)
        
        if current_slide:
            slides.append(current_slide)
        
        if len(slides) < 2:
            return None
        
        return slides

def create_presentation(prompt: str, context: str, template_name: str | None = None) -> str:
    """
    Generate a PowerPoint presentation from prompt and context.
    
    Args:
        prompt: User's prompt for the presentation
        context: Additional context
        template_name: Optional template name (or random if None)
        
    Returns:
        Path to the generated .pptx file
    """
    system_message = """You are an expert presentation designer. Create a structured PowerPoint presentation.
    Return your response as a JSON object with this exact structure:
    {
        "slides": [
            {
                "title": "Slide Title",
                "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"]
            }
        ]
    }
    
    Rules:
    - Create 5-8 slides total
    - First slide should be a title slide with the main topic
    - Each content slide should have 3-5 bullet points
    - Keep bullet points concise (max 10-15 words each)
    - Last slide should be a conclusion or summary
    - Use clear, professional language
    """
    
    full_prompt = f"{prompt}\n\nContext: {context}"
    
    slides_data = None
    
    try:
        ai_response = generate_content(
            prompt=full_prompt,
            system_message=system_message,
            max_completion_tokens=4096
        )
        slides_data = parse_ai_slides(ai_response)
    except Exception as e:
        print(f"AI generation error: {e}")
        slides_data = None
    
    if not slides_data or len(slides_data) < 2:
        slides_data = create_fallback_slides(prompt, context)
    
    template = get_template(template_name)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for idx, slide_data in enumerate(slides_data):
        if idx == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get("title", "Presentation")
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = template["title_color"]
            title_para.alignment = PP_ALIGN.CENTER
            
            if slide_data.get("content"):
                subtitle_box = slide.shapes.add_textbox(
                    Inches(1), Inches(4.5), Inches(8), Inches(1)
                )
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = " | ".join(slide_data["content"][:2])
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(20)
                subtitle_para.font.color.rgb = template["text_color"]
                subtitle_para.alignment = PP_ALIGN.CENTER
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get("title", f"Slide {idx + 1}")
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = template["title_color"]
            
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.8), Inches(8.4), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for bullet in slide_data.get("content", []):
                p = content_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = template["text_color"]
                p.space_before = Pt(12)
    
    os.makedirs("generated_ppts", exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"generated_ppts/presentation_{timestamp}.pptx"
    
    prs.save(filename)
    return filename
